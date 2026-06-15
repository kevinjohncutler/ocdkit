"""Tests for ocdkit.io.pptx.figs_to_deck — mpl Figures, SvgFigure, and the
LiveFigure (duck-typed ``to_svg()``) handle, plus per-slide title-pull."""
import zipfile

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
import pytest  # noqa: E402
from pptx import Presentation  # noqa: E402

from ocdkit.io import LiveFigure, figs_to_deck  # noqa: E402


def _mpl_fig(label="x"):
    f = plt.figure(figsize=(4, 3))
    ax = f.add_subplot(111)
    ax.plot([0, 1, 2], [0, 1, 0])
    ax.set_title(label)
    return f


def _live(label):
    # to_svg() returns an mpl Figure, which _unwrap_figs handles via the
    # MplFigure branch — exercises the duck-typed handle path without wgpu.
    return LiveFigure(render=lambda *, as_svg=True, out=None: _mpl_fig(label), title=label)


class TestFigsToDeck:
    def test_mpl_figures(self, tmp_path):
        out = tmp_path / "deck.pptx"
        figs_to_deck([_mpl_fig("a"), _mpl_fig("b")], str(out),
                     verbose=False, save_debug=False, save_crops=False)
        assert out.exists() and out.stat().st_size > 0
        assert len(Presentation(str(out)).slides) == 2

    def test_livefigure_handles(self, tmp_path):
        out = tmp_path / "live.pptx"
        figs = [_live("Slide A"), _live("Slide B"), _live("Slide C")]
        figs_to_deck(figs, str(out), verbose=False, save_debug=False, save_crops=False)
        assert out.exists()
        assert len(Presentation(str(out)).slides) == 3

    def test_titles_pulled_from_handles(self, tmp_path):
        out = tmp_path / "titled.pptx"
        # no titles= kwarg -> figs_to_deck pulls each slide title from the handle
        figs_to_deck([_live("My Unique Slide Title")], str(out),
                     verbose=False, save_debug=False, save_crops=False)
        with zipfile.ZipFile(out) as z:
            blob = b"".join(z.read(n) for n in z.namelist() if n.endswith(".xml"))
        assert b"My Unique Slide Title" in blob

    def test_rejects_unsupported_input(self, tmp_path):
        with pytest.raises(TypeError):
            figs_to_deck([12345], str(tmp_path / "bad.pptx"),
                         verbose=False, save_debug=False, save_crops=False)
