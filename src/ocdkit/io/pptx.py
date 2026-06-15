
#need an option on spectral color plot to recolor IR


import re as _re

_SVG_ROOT_RE = _re.compile(rb"(<svg\b[^>]*?)>", _re.DOTALL)
_SVG_WIDTH_ATTR_RE = _re.compile(rb'\s+width="[^"]*"')
_SVG_HEIGHT_ATTR_RE = _re.compile(rb'\s+height="[^"]*"')


def _set_svg_root_dimensions(svg_bytes: bytes,
                              *, width_in: float, height_in: float) -> bytes:
    """Rewrite the ``<svg>`` root's ``width`` / ``height`` to fixed inches.

    Sets ``width="<W>in" height="<H>in"`` on the FIRST ``<svg>`` tag
    (the root) so PowerPoint's SVG renderer + convert-to-shapes use
    these as the canvas size — making SVG user units map 1:1 to the
    intended on-slide geometry.  Any pre-existing ``width`` / ``height``
    attributes on the root are removed first to avoid duplicates.
    """
    m = _SVG_ROOT_RE.search(svg_bytes)
    if not m:
        return svg_bytes
    open_tag = m.group(1)
    open_tag = _SVG_WIDTH_ATTR_RE.sub(b"", open_tag, count=1)
    open_tag = _SVG_HEIGHT_ATTR_RE.sub(b"", open_tag, count=1)
    new_attrs = f' width="{width_in:.4f}in" height="{height_in:.4f}in"'.encode()
    open_tag = open_tag + new_attrs
    return svg_bytes[:m.start()] + open_tag + b">" + svg_bytes[m.end():]


def _unwrap_figs(figs, caller):
    """Normalize an iterable of figs to a list of matplotlib Figures.

    Accepts:
      * ``matplotlib.figure.Figure`` — passed through unchanged.
      * :class:`ocdkit.io.SvgFigure` (SVG-backed) — rasterized via
        ``render_to_image`` and wrapped in a minimal mpl Figure
        containing the resulting image (so the downstream precise-
        export pipeline, which is matplotlib-native, keeps working
        unchanged).  The original SVG payload is stashed on the
        bridge figure so the PPTX writer can also embed the vector
        SVG alongside the rasterized PNG.
      * Any lazy handle exposing ``to_svg()`` (e.g.
        :class:`ocdkit.io.LiveFigure`) — rendered to an ``SvgFigure``
        on demand, then treated as above.

    Anything else raises a TypeError naming the offending type.
    """
    from matplotlib.figure import Figure as _MplFigure
    from ocdkit.io import SvgFigure as _SvgFigure

    out = []
    for i, fig in enumerate(figs):
        # Lazy live handles (LiveFigure and any object exposing ``to_svg()``)
        # render to an SvgFigure on demand, then flow through the SVG path.
        if not isinstance(fig, (_MplFigure, _SvgFigure)) and callable(getattr(fig, "to_svg", None)):
            fig = fig.to_svg()
        if isinstance(fig, _MplFigure):
            out.append(fig)
        elif isinstance(fig, _SvgFigure):
            out.append(_svgfigure_to_mpl(fig))
        else:
            hint = ""
            if isinstance(fig, str):
                hint = (" (got a str — a plotting call returned a raw SVG "
                        "string instead of a figure object; pass the "
                        "Figure / SvgFigure / LiveFigure handle directly)")
            raise TypeError(
                f"{caller}: figs[{i}] must be a matplotlib.figure.Figure, "
                f"ocdkit.io.SvgFigure, or a LiveFigure/to_svg() handle, "
                f"got {type(fig).__name__}{hint}"
            )
    return out


_SVG_BRIDGE_RENDER_DPI = 200  # quality of the SVG → raster handoff
_SVG_BRIDGE_MAX_PIXELS = 50_000_000  # cap to avoid Pillow's bomb guard


def _svgfigure_to_mpl(svg_fig):
    """Rasterize an SvgFigure and wrap the result in a minimal mpl Figure.

    The downstream PPTX pipeline expects a matplotlib Figure (it crops,
    snaps, embeds via the matplotlib axes geometry).  Wrapping the
    rasterized SVG in a one-axes mpl Figure lets the existing code path
    handle SVG sources without rewriting the exporter.

    The intermediate mpl Figure is sized so that ``savefig`` at *any*
    downstream DPI reproduces the rasterized SVG without re-upscaling
    — figsize is set such that ``pixels_out = arr.shape`` regardless
    of the dpi argument.  Caps the intermediate at
    ``_SVG_BRIDGE_MAX_PIXELS`` to stay below Pillow's bomb guard.

    The original SVG payload is stashed on the bridge figure as
    ``_ocdkit_svg_source`` so the PPTX writer can embed the vector
    SVG alongside the rasterized PNG fallback.
    """
    w_in, h_in = svg_fig.intrinsic_size_in
    # Choose a render DPI that respects the pixel cap.
    target_px = w_in * h_in * _SVG_BRIDGE_RENDER_DPI * _SVG_BRIDGE_RENDER_DPI
    if target_px > _SVG_BRIDGE_MAX_PIXELS:
        import math
        render_dpi = int(math.sqrt(_SVG_BRIDGE_MAX_PIXELS / (w_in * h_in)))
    else:
        render_dpi = _SVG_BRIDGE_RENDER_DPI
    arr = svg_fig.render_to_image(dpi=render_dpi)
    bridge = _array_to_mpl_figure(arr, render_dpi)
    # Stash the PPTX-embeddable SVG (JXL → PNG, no fixed width/height) so
    # the PPTX exporter can embed something PowerPoint actually
    # understands AND so convert-to-shapes scales fonts to the slide
    # placement instead of the SVG's native pixel canvas.
    bridge._ocdkit_svg_source = svg_fig._pptx_embeddable_svg()
    return bridge


def _attach_svg_to_pptx(pptx_path, slide_idx_to_svg, *, rewrite_dims=True):
    """Post-process a PPTX to attach SVG variants to specific slides.

    PowerPoint 2016+ supports an ``asvg:svgBlip`` extension on a
    picture's ``a:blip`` — when present, PowerPoint renders the SVG;
    older clients fall back to the embedded PNG.  We keep BOTH so
    fidelity is gained where supported and nothing regresses where
    it isn't.

    Parameters
    ----------
    pptx_path : Path-like
        The .pptx file produced by python-pptx (will be modified in place).
    slide_idx_to_svg : dict[int, str | bytes]
        ``{1-based slide index: svg payload}``.  For each slide
        present in this dict, *every* ``<p:pic>`` in that slide gets
        an ``asvg:svgBlip`` pointing at the SVG.  Use this only for
        slides whose pictures should all share the same SVG vector
        source (the SvgFigure → bridge path produces a single
        whole-canvas picture per slide, so this fits naturally).
    """
    import re
    import shutil
    import zipfile
    from pathlib import Path
    from lxml import etree

    if not slide_idx_to_svg:
        return

    pptx_path = Path(pptx_path)
    work = pptx_path.with_suffix(".tmp.pptx")

    NS = {
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "asvg": "http://schemas.microsoft.com/office/drawing/2016/SVG/main",
        "ct": "http://schemas.openxmlformats.org/package/2006/content-types",
        "rels": "http://schemas.openxmlformats.org/package/2006/relationships",
    }
    BLIP_SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
    EMU_PER_INCH = 914400.0

    # Normalize SVG bytes.
    svg_bytes_by_slide = {
        idx: (s.encode("utf-8") if isinstance(s, str) else s)
        for idx, s in slide_idx_to_svg.items()
    }

    with zipfile.ZipFile(pptx_path, "r") as zin:
        names = zin.namelist()
        data_by_name = {n: zin.read(n) for n in names}

    # Pick next free image index across the package.
    used_indices = set()
    for n in names:
        m = re.match(r"^ppt/media/image(\d+)\.\w+$", n)
        if m:
            used_indices.add(int(m.group(1)))
    next_img_idx = max(used_indices, default=0) + 1

    # Plan: for each slide, allocate an image filename + relationship id.
    # The rid is per-slide (rels are per-slide), the image filename is global.
    plan = {}  # slide_idx -> dict(svg_name, svg_bytes, rid)
    for slide_idx in sorted(svg_bytes_by_slide):
        plan[slide_idx] = {
            "svg_name": f"image{next_img_idx}.svg",
            "svg_bytes": svg_bytes_by_slide[slide_idx],
            "rid": None,  # filled in when we patch the rels
        }
        next_img_idx += 1

    # Patch each affected slide-rels XML and slide XML.
    for slide_idx, info in plan.items():
        rels_key = f"ppt/slides/_rels/slide{slide_idx}.xml.rels"
        slide_key = f"ppt/slides/slide{slide_idx}.xml"
        if rels_key not in data_by_name or slide_key not in data_by_name:
            continue

        # 1. Add relationship to the SVG media file.
        rels_root = etree.fromstring(data_by_name[rels_key])
        existing_rids = [r.get("Id") for r in rels_root]
        next_rid = max(
            (int(r[3:]) for r in existing_rids
             if r.startswith("rId") and r[3:].isdigit()),
            default=0,
        ) + 1
        rid = f"rId{next_rid}"
        info["rid"] = rid
        new_rel = etree.SubElement(rels_root, f"{{{NS['rels']}}}Relationship")
        new_rel.set("Id", rid)
        new_rel.set("Type",
                     "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image")
        new_rel.set("Target", f"../media/{info['svg_name']}")
        data_by_name[rels_key] = etree.tostring(
            rels_root, xml_declaration=True, encoding="UTF-8", standalone=True)

        # 2. Patch every <a:blip> in this slide to add asvg:svgBlip.
        slide_root = etree.fromstring(data_by_name[slide_key])
        for blip in slide_root.xpath(".//a:blip", namespaces=NS):
            if blip.xpath("./a:extLst/a:ext/asvg:svgBlip", namespaces=NS):
                continue  # already has one
            ext_lst = etree.SubElement(blip, f"{{{NS['a']}}}extLst")
            ext = etree.SubElement(ext_lst, f"{{{NS['a']}}}ext", uri=BLIP_SVG_EXT_URI)
            svg_blip = etree.SubElement(ext, f"{{{NS['asvg']}}}svgBlip")
            svg_blip.set(f"{{{NS['r']}}}embed", rid)
        data_by_name[slide_key] = etree.tostring(
            slide_root, xml_declaration=True, encoding="UTF-8", standalone=True)

        # 3. Rewrite the SVG payload's root <svg width/height> to match
        #    the on-slide picture box.  Without this, PowerPoint
        #    interprets the SVG's native dimensions (e.g. 600px = 6.25 in)
        #    as the canvas when converting picture → editable shapes, so
        #    font-size="5" becomes 5px ≈ 3.75pt OR (if width/height are
        #    missing) it falls back to weird defaults that make text
        #    explode.  Pinning width/height to the picture-box EMU value
        #    lets shape-convert use the right viewBox→inch ratio and font
        #    sizes (which live in SVG user units) come out at the size we
        #    laid them out for.
        pics = slide_root.xpath(".//p:pic", namespaces=NS)
        if pics and rewrite_dims:
            ext_el = pics[0].xpath(".//p:spPr/a:xfrm/a:ext", namespaces=NS)
            if ext_el:
                try:
                    cx = int(ext_el[0].get("cx", "0"))
                    cy = int(ext_el[0].get("cy", "0"))
                    if cx > 0 and cy > 0:
                        info["svg_bytes"] = _set_svg_root_dimensions(
                            info["svg_bytes"],
                            width_in=cx / EMU_PER_INCH,
                            height_in=cy / EMU_PER_INCH,
                        )
                except (TypeError, ValueError):
                    pass  # leave SVG dims as-is if EMU parse fails

    # 3. Ensure [Content_Types].xml has a Default for svg.
    ct_key = "[Content_Types].xml"
    if ct_key in data_by_name and b"image/svg+xml" not in data_by_name[ct_key]:
        ct_root = etree.fromstring(data_by_name[ct_key])
        new_default = etree.SubElement(ct_root, f"{{{NS['ct']}}}Default")
        new_default.set("Extension", "svg")
        new_default.set("ContentType", "image/svg+xml")
        data_by_name[ct_key] = etree.tostring(
            ct_root, xml_declaration=True, encoding="UTF-8", standalone=True)

    # Rewrite the zip with patched files + the new SVG media parts.
    with zipfile.ZipFile(work, "w", zipfile.ZIP_DEFLATED) as zout:
        for name in names:
            zout.writestr(name, data_by_name[name])
        for slide_idx, info in plan.items():
            if info["rid"] is None:
                continue  # slide wasn't in the package
            zout.writestr(f"ppt/media/{info['svg_name']}", info["svg_bytes"])

    shutil.move(str(work), str(pptx_path))


# ─────────────────────── convert-to-shapes fidelity ──────────────────────────
# NOTE: the PowerPoint "Convert to Shapes" SVG fidelity fixes moved to
# ocdkit.io.figure.apply_convert_to_shapes_fixes (next to the svgBlip
# embedding they correct); imported at the call site below.


def _array_to_mpl_figure(arr, render_dpi):
    """Build a minimal mpl Figure that reproduces ``arr`` 1:1 on savefig.

    figsize is set so each pixel of ``arr`` maps to one output pixel at
    ``render_dpi``.  Downstream ``fig.savefig(..., dpi=X)`` will scale
    the canvas by ``X / render_dpi`` — bounded, not the catastrophic
    inches-times-final-dpi blowup we'd get from a figsize-in-inches
    bridge.
    """
    import matplotlib.pyplot as plt
    h, w = arr.shape[:2]
    fig, ax = plt.subplots(figsize=(w / render_dpi, h / render_dpi),
                            dpi=render_dpi)
    ax.imshow(arr)
    ax.set_position([0, 0, 1, 1])
    ax.axis("off")
    return fig


def figs_to_deck(
    figs,
    pptx_path,
    dpi=200,
    fontcolor=None,
    axcolor=None,
    # *,
    title: str | None = None,
    titles: list[str | None] | None = None,
    show_slide_number: bool = True,
    background_rgb: tuple[int, int, int] = (0, 0, 0),
    template_path: str | None = None,
    layout_name: str | None = None,
    slide_padding_in: float = 0.0,
    title_font_pt: int = 28,
    title_font_name: str | None = None,
    title_color_rgb: tuple[int, int, int] = (255, 255, 255),
    title_height_in: float = 0.7,
    title_align: str = "center",
    footer_reserved_in: float = 0.35,
    slide_number_pos: str = "right",
    slide_number_font_pt: int = 12,
    slide_number_color_rgb: tuple[int, int, int] = (255, 255, 255),
    slide_number_margin_in: float = 0.2,
    slide_number_font_name: str | None = None,
    content_box_inches: tuple[float, float, float, float] | None = None,
    verbose: bool = False,
    save_crops: bool = False,
    crops_dir: str | None = None,
    snap_to_pixel: bool = True,
    scale_mode: str = "fit",
    write_debug: bool = True,
    save_debug: bool = False,
):
    """
    Build a PPTX deck by delegating to figs_to_deck_precise_combined to ensure
    identical behavior to the demo path. Applies optional color overrides to
    figures before export, then constructs a per-slide options list.

    Per-slide titles default to each input handle's ``.title`` attribute
    (e.g. a :class:`ocdkit.io.LiveFigure`) when ``titles=`` is not supplied.
    """
    figs = list(figs)
    if titles is None:
        _handle_titles = [getattr(f, "title", None) for f in figs]
        if any(t is not None for t in _handle_titles):
            titles = _handle_titles
    # Pre-apply the color scheme on Figure-abstraction inputs BEFORE
    # rasterizing.  For SvgFigure the recolor mutates the underlying
    # SVG, so the colors carry into both the PNG rasterization (via
    # the bridge) and the natively-embedded SVG variant.  Also drop
    # the SVG's own opaque background so the slide background shows
    # through naturally — otherwise the SVG renders as a coloured card
    # sitting on top of whatever ``background_rgb`` was set.
    if fontcolor is not None or axcolor is not None:
        from ocdkit.io import SvgFigure as _SvgFigure
        for fig in figs:
            if isinstance(fig, _SvgFigure):
                fig.apply_color_scheme(font=fontcolor, axes=axcolor)
                fig.set_facecolor("none")
    # Normalize: accept mpl Figures, SvgFigure, or LiveFigure handles; SVG-backed
    # ones get rasterized so the legacy precise-export path (which is
    # matplotlib-native) can consume them unchanged.
    figs = _unwrap_figs(figs, "figs_to_deck")
    # Optional color overrides before cropping/rasterization
    if fontcolor is not None or axcolor is not None:
        from matplotlib.text import Text
        for fig in figs:
            for ax in fig.get_axes():
                if fontcolor is not None:
                    ax.title.set_color(fontcolor)
                    if ax.xaxis is not None:
                        ax.xaxis.label.set_color(fontcolor)
                    if ax.yaxis is not None:
                        ax.yaxis.label.set_color(fontcolor)
                    ax.tick_params(colors=fontcolor)
                    for txt in getattr(ax, "texts", []):
                        try:
                            txt.set_color(fontcolor)
                        except Exception:
                            pass
                if axcolor is not None:
                    for sp in getattr(ax, "spines", {}).values():
                        try:
                            sp.set_color(axcolor)
                        except Exception:
                            pass
                    ax.tick_params(color=axcolor)
            if fontcolor is not None:
                for txt in fig.findobj(Text):
                    try:
                        txt.set_color(fontcolor)
                    except Exception:
                        pass

    # Prepare slides list with per-slide opts (identical to combined path)
    # Default crops_dir to debug/ if not provided so artifacts are easy to find.
    if not save_debug:
        if crops_dir is not None and verbose:
            try:
                print("Ignoring crops_dir because save_debug=False")
            except Exception:
                pass
        crops_dir = None
    elif crops_dir is None:
        try:
            from . import fig_export as FE
            crops_dir = str(FE._ensure_debug_dir())
        except Exception:
            crops_dir = None
    # ``titles`` lets the caller annotate each slide independently while
    # reusing the existing ``title_font_pt`` / ``title_color_rgb`` /
    # ``title_align`` / ``title_height_in`` / ``title_font_name`` controls
    # (no per-slide text-box management on the caller side). Missing
    # entries (or a shorter list) fall back to the single ``title``.
    if titles is not None and len(titles) != len(figs):
        titles = list(titles) + [None] * max(0, len(figs) - len(titles))
    slides = []
    for i, fig in enumerate(figs):
        per_slide_title = title
        if titles is not None and titles[i] is not None:
            per_slide_title = titles[i]
        opts = dict(
            dpi=dpi,
            background_rgb=background_rgb,
            layout_name=layout_name,
            title=per_slide_title,
            title_font_pt=title_font_pt,
            title_color_rgb=title_color_rgb,
            title_height_in=title_height_in,
            title_align=title_align,
            title_font_name=title_font_name,
            content_box_inches=content_box_inches,
            slide_padding_in=slide_padding_in,
            footer_reserved_in=footer_reserved_in,
            show_slide_number=show_slide_number,
            slide_number_pos=slide_number_pos,
            slide_number_font_pt=slide_number_font_pt,
            slide_number_color_rgb=slide_number_color_rgb,
            slide_number_margin_in=slide_number_margin_in,
            slide_number_font_name=slide_number_font_name,
            save_crops=save_crops,
            crops_dir=crops_dir,
            snap_to_pixel=snap_to_pixel,
            scale_mode=scale_mode,
        )
        slides.append((fig, opts))

    # Delegate to the combined precise exporter for identical behavior
    return figs_to_deck_precise_combined(
        slides,
        pptx_path,
        template_path=template_path,
        verbose=verbose,
        save_crops=save_crops,
        save_debug=save_debug,
        crops_dir=crops_dir,
        snap_to_pixel=snap_to_pixel,
        write_debug=write_debug,
    )


def figs_to_deck_precise(
    figs,
    pptx_path,
    dpi=200,
    background_rgb=(0, 0, 0),
    template_path: str | None = None,
    layout_name: str | None = None,
    title: str | None = None,
    title_font_pt: int = 28,
    title_color_rgb: tuple[int, int, int] = (255, 255, 255),
    title_height_in: float = 0.7,
    title_align: str = "center",
    title_font_name: str | None = None,
    content_box_inches: tuple[float, float, float, float] | None = None,
    slide_padding_in: float | tuple[float, float, float, float] = 0.0,
    footer_reserved_in: float = 0.35,
    show_slide_number: bool = True,
    slide_number_pos: str = "right",
    slide_number_font_pt: int = 12,
    slide_number_color_rgb: tuple[int, int, int] = (255, 255, 255),
    slide_number_margin_in: float = 0.2,
    slide_number_font_name: str | None = None,
    slide_number_mode: str = "template",
    verbose: bool = False,
    save_crops: bool = False,
    crops_dir: str | None = None,
    snap_to_pixel: bool = True,
    scale_mode: str = "fit",
    write_debug: bool = True,
    save_debug: bool = False,
):
    """
    Export figures as multiple precisely placed subplot images.

    Delegates to figs_to_deck_precise_combined so behavior stays aligned with
    the higher-level path used by figs_to_deck.
    """
    figs = _unwrap_figs(figs, "figs_to_deck_precise")
    slides = []
    for fig in figs:
        opts = dict(
            dpi=dpi,
            background_rgb=background_rgb,
            layout_name=layout_name,
            title=title,
            title_font_pt=title_font_pt,
            title_color_rgb=title_color_rgb,
            title_height_in=title_height_in,
            title_align=title_align,
            title_font_name=title_font_name,
            content_box_inches=content_box_inches,
            slide_padding_in=slide_padding_in,
            footer_reserved_in=footer_reserved_in,
            show_slide_number=show_slide_number,
            slide_number_pos=slide_number_pos,
            slide_number_font_pt=slide_number_font_pt,
            slide_number_color_rgb=slide_number_color_rgb,
            slide_number_margin_in=slide_number_margin_in,
            slide_number_font_name=slide_number_font_name,
            save_crops=save_crops,
            crops_dir=crops_dir,
            snap_to_pixel=snap_to_pixel,
            scale_mode=scale_mode,
        )
        slides.append((fig, opts))

    return figs_to_deck_precise_combined(
        slides,
        pptx_path,
        template_path=template_path,
        verbose=verbose,
        save_crops=save_crops,
        save_debug=save_debug,
        crops_dir=crops_dir,
        snap_to_pixel=snap_to_pixel,
        write_debug=write_debug,
    )

def _add_precise_slide(
    prs,
    fig,
    *,
    dpi=200,
    background_rgb=(0, 0, 0),
    layout_name: str | None = None,
    title: str | None = None,
    title_font_pt: int = 28,
    title_color_rgb: tuple[int, int, int] = (255, 255, 255),
    title_height_in: float = 0.7,
    title_align: str = "center",
    title_font_name: str | None = None,
    content_box_inches: tuple[float, float, float, float] | None = None,
    slide_padding_in: float | tuple[float, float, float, float] = 0.0,
    footer_reserved_in: float = 0.35,
    show_slide_number: bool = True,
    slide_number_pos: str = "right",
    slide_number_font_pt: int = 12,
    slide_number_color_rgb: tuple[int, int, int] = (255, 255, 255),
    slide_number_margin_in: float = 0.2,
    slide_number_font_name: str | None = None,
    save_crops: bool = False,
    crops_dir: str | None = None,
    crops_prefix: str | None = None,
    snap_to_pixel: bool = True,
    scale_mode: str = "fit",
    save_debug: bool = False,
    debug_dir=None,
):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from . import fig_export as FE
    import io as _io

    slide_w, slide_h = prs.slide_width, prs.slide_height
    
    def _choose_layout(prs_):
        if layout_name:
            for layout in prs_.slide_layouts:
                try:
                    if layout.name and layout.name.strip().lower() == layout_name.strip().lower():
                        return layout
                except Exception:
                    pass
        # Prefer any layout named like 'Blank'
        for layout in prs_.slide_layouts:
            try:
                if layout.name and "blank" in layout.name.strip().lower():
                    return layout
            except Exception:
                pass
        # Fallback: index 6 then 0
        try:
            return prs_.slide_layouts[6]
        except Exception:
            return prs_.slide_layouts[0]

    # Normalize padding
    if isinstance(slide_padding_in, (int, float)):
        pad_l_in = pad_t_in = pad_r_in = pad_b_in = float(slide_padding_in)
    else:
        pad_l_in, pad_t_in, pad_r_in, pad_b_in = slide_padding_in
    pad_l = Inches(pad_l_in); pad_t = Inches(pad_t_in); pad_r = Inches(pad_r_in); pad_b = Inches(pad_b_in)

    # Collect axes crops and figure canvas
    if save_debug:
        if debug_dir is None:
            debug_dir = FE._ensure_debug_dir()
        bbox_log = debug_dir / "pptx_axes_bboxes.txt"
        crops_dir = crops_dir or str(debug_dir)
    else:
        debug_dir = None
        bbox_log = None
        save_crops = False
        crops_dir = None
    # Export per-axes crops using fig_export's union-bbox approach
    crops = FE.export_axes_to_buffers(
        fig,
        dpi=dpi,
        pad_inches=0.0,
        log_path=bbox_log,
        include_legend=True,
        pad_px=0,
        hide_figure_text=True,
        include_tick_marks=True,
        tick_label_clip_margin_px=12,
        snap_to_pixel=snap_to_pixel,
    )
    # Ensure figure pixel size corresponds to the same DPI used for cropping
    _orig_dpi = fig.dpi
    full_fig_png = None
    try:
        fig.set_dpi(dpi)
        canvas = FE.FigureCanvas(fig); canvas.draw()
        fig_w_px, fig_h_px = canvas.get_width_height()
        try:
            buf_full = _io.BytesIO()
            fig.savefig(buf_full, format="png", dpi=dpi, facecolor=fig.get_facecolor())
            full_fig_png = buf_full.getvalue()
        except Exception:
            full_fig_png = None
    finally:
        try:
            fig.set_dpi(_orig_dpi)
        except Exception:
            pass

    try:
        fig_text_crops = FE.export_figure_texts_to_buffers(fig, dpi=dpi, pad_inches=0.0)
    except Exception:
        fig_text_crops = []

    # New blank slide and background
    chosen_layout = _choose_layout(prs)
    slide = prs.slides.add_slide(chosen_layout)
    fill = slide.background.fill
    if background_rgb is not None:
        fill.solid(); fill.fore_color.rgb = RGBColor(*background_rgb)

    # Compute content box
    if content_box_inches is not None:
        left_in, top_in, width_in, height_in = content_box_inches
        content_left = Inches(left_in); content_top = Inches(top_in)
        content_w = Inches(width_in);   content_h = Inches(height_in)
    else:
        content_left = pad_l; content_top = pad_t
        content_w = slide_w - pad_l - pad_r
        content_h = slide_h - pad_t - pad_b
        if title:
            lines = [title] if "\n" not in title else title.splitlines()
            base_th = Inches(title_height_in)
            approx_line_h = int(Pt(title_font_pt) * 1.25)
            approx_total_h = approx_line_h * max(1, len(lines)) + int(Pt(title_font_pt) * 0.25)
            th = max(base_th, approx_total_h)
            title_box = slide.shapes.add_textbox(content_left, content_top, content_w, th)
            tf = title_box.text_frame; tf.clear()
            for idx, line in enumerate(lines):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                run = p.add_run(); run.text = line
                run.font.size = Pt(title_font_pt); run.font.bold = True
                if title_font_name: run.font.name = title_font_name
                run.font.color.rgb = RGBColor(*title_color_rgb)
                _align = str(title_align).lower() if title_align is not None else "center"
                if _align == "left":
                    p.alignment = PP_ALIGN.LEFT
                elif _align == "right":
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.CENTER
            content_top = content_top + th
            content_h = slide_h - content_top - pad_b

    if footer_reserved_in and footer_reserved_in > 0:
        reserve = Inches(footer_reserved_in)
        if content_h > reserve:
            content_h -= reserve

    # Determine effective bounds (in figure pixels) that must fit on the slide.
    placement_pad_px = 0
    combined_extents = []

    def _append_extent(x0: float, y0: float, x1: float, y1: float) -> None:
        combined_extents.append((float(x0), float(y0), float(x1), float(y1)))

    for item in crops:
        bb = item["bbox_disp"]
        _append_extent(bb.x0 - placement_pad_px, bb.y0 - placement_pad_px,
                       bb.x1 + placement_pad_px, bb.y1 + placement_pad_px)

    for item in fig_text_crops:
        bb = item["bbox_disp"]
        _append_extent(bb.x0, bb.y0, bb.x1, bb.y1)

    if not combined_extents:
        _append_extent(0.0, 0.0, float(fig_w_px), float(fig_h_px))

    full_x0 = min(ext[0] for ext in combined_extents)
    full_y0 = min(ext[1] for ext in combined_extents)
    full_x1 = max(ext[2] for ext in combined_extents)
    full_y1 = max(ext[3] for ext in combined_extents)

    total_width_px = max(full_x1 - full_x0, 1.0)
    total_height_px = max(full_y1 - full_y0, 1.0)

    if content_w <= 0 or content_h <= 0:
        raise ValueError("Computed content area is non-positive; check padding/layout inputs")

    mode = (scale_mode or "fit").lower()
    if mode not in {"fit", "native", "fill", "auto"}:
        mode = "fit"

    fig_w_in = float(fig.get_figwidth() or 0.0)
    fig_h_in = float(fig.get_figheight() or 0.0)
    desired_w = Inches(fig_w_in) if fig_w_in > 0 else None
    desired_h = Inches(fig_h_in) if fig_h_in > 0 else None
    use_full_frame_native = bool(full_fig_png) and mode == "native" and desired_w and desired_h

    def _fit_scale() -> tuple[float, int, int, float, float]:
        sc_w = content_w / total_width_px
        sc_h = content_h / total_height_px
        sc = min(sc_w, sc_h) if mode == "fit" else max(sc_w, sc_h)
        actual_w = sc * total_width_px
        actual_h = sc * total_height_px
        off_left = content_left + int(round((content_w - actual_w) / 2))
        off_top = content_top + int(round((content_h - actual_h) / 2))
        return sc, off_left, off_top, actual_w, actual_h

    if mode == "native":
        if desired_w and desired_h:
            scale_w = desired_w / total_width_px
            scale_h = desired_h / total_height_px
            scale = min(scale_w, scale_h)
            actual_w = scale * total_width_px
            actual_h = scale * total_height_px
            offset_left = content_left + int(round((content_w - actual_w) / 2))
            offset_top = content_top + int(round((content_h - actual_h) / 2))
            if use_full_frame_native:
                actual_w = float(desired_w)
                actual_h = float(desired_h)
                scale = float(desired_w) / max(total_width_px, 1.0)
                offset_left = content_left + int(round((content_w - actual_w) / 2))
                offset_top = content_top + int(round((content_h - actual_h) / 2))
        else:
            scale, offset_left, offset_top, actual_w, actual_h = _fit_scale()
    else:
        scale, offset_left, offset_top, actual_w, actual_h = _fit_scale()

    # Place subplots
    base_fig_info = None
    if use_full_frame_native:
        try:
            base_left = int(round(offset_left))
            base_top = int(round(offset_top))
            base_width = int(round(actual_w))
            base_height = int(round(actual_h))
            buf_full = _io.BytesIO(full_fig_png); buf_full.seek(0)
            slide.shapes.add_picture(buf_full, base_left, base_top, width=base_width, height=base_height)
            base_fig_info = (base_left, base_top, base_width, base_height)
        except Exception:
            base_fig_info = None

    placements = []
    if not use_full_frame_native:
        for item in crops:
            bb = item["bbox_disp"]
            bb_x0 = bb.x0 - placement_pad_px
            bb_y0 = bb.y0 - placement_pad_px
            bb_x1 = bb.x1 + placement_pad_px
            bb_y1 = bb.y1 + placement_pad_px
            w_px = bb_x1 - bb_x0
            h_px = bb_y1 - bb_y0
            left_emu = int(round(offset_left + scale * (bb_x0 - full_x0)))
            top_emu = int(round(offset_top + scale * (full_y1 - bb_y1)))
            width_emu = int(round(scale * w_px))
            height_emu = int(round(scale * h_px))

            buf = _io.BytesIO(item["png"]); buf.seek(0)
            slide.shapes.add_picture(buf, left_emu, top_emu, width=width_emu, height=height_emu)
            placements.append((item["index"], left_emu, top_emu, width_emu, height_emu,
                               int(round(bb.width)), int(round(bb.height))))

            # Optionally save the exact PNG crop for inspection
            if save_crops:
                try:
                    from pathlib import Path as _P
                    odir = _P(crops_dir) if crops_dir else _P('.')
                    odir.mkdir(parents=True, exist_ok=True)
                    prefix = crops_prefix or 'crops'
                    fname = f"{prefix}_a{item['index']:02d}.png"
                    with open(odir / fname, 'wb') as fh:
                        fh.write(buf.getvalue())
                except Exception:
                    pass

    # Placement log
    if debug_dir is not None:
        place_path = debug_dir / "pptx_axes_placements.txt"
        with open(place_path, "a", encoding="utf-8") as fh:
            print("", file=fh)
            print(f"[Slide] layout='{slide.slide_layout.name}' background_rgb={background_rgb}", file=fh)
            # No placeholder logging; we inject slide numbers directly.
            # Log shapes count (helps diagnose background or overlay issues)
            try:
                total = len(slide.shapes)
                pics = sum(1 for s in slide.shapes if getattr(s, "shape_type", None) and str(s.shape_type).endswith("PICTURE"))
                tboxes = sum(1 for s in slide.shapes if getattr(s, "has_text_frame", False))
                print(f"Shapes: total={total} pictures={pics} textboxes={tboxes}", file=fh)
            except Exception:
                pass
            print(f"Figure canvas: {fig_w_px} x {fig_h_px} px", file=fh)
            print(f"Effective bounds px: ({full_x0:.1f}, {full_y0:.1f})→({full_x1:.1f}, {full_y1:.1f})", file=fh)
            print(f"Content box (EMU): left={content_left} top={content_top} width={content_w} height={content_h}", file=fh)
            print(f"Scale (EMU/px): {scale}", file=fh)
            if base_fig_info:
                b_left, b_top, b_w, b_h = base_fig_info
                print(f"Base figure (native): left={b_left} top={b_top} width={b_w} height={b_h}", file=fh)
            print(f"Effective size px: {total_width_px:.1f} x {total_height_px:.1f}", file=fh)
            if footer_reserved_in:
                print(f"Footer reserved (EMU): {Inches(footer_reserved_in)}", file=fh)
            for (i, L, T, W, H, wp, hp) in placements:
                print(f"axes {i:02d}: left={L} top={T} width={W} height={H} (px: {wp}x{hp}) pad_px={placement_pad_px}", file=fh)
            if fig_text_crops and not use_full_frame_native:
                print(f"Figure-level texts: {len(fig_text_crops)}", file=fh)

    # Add figure-level texts as raster images for precise placement
    if not use_full_frame_native:
        for item in fig_text_crops:
            try:
                bb = item["bbox_disp"]
                w_px = bb.x1 - bb.x0
                h_px = bb.y1 - bb.y0
                if w_px <= 0 or h_px <= 0:
                    continue
                left_emu = int(round(offset_left + scale * (bb.x0 - full_x0)))
                top_emu = int(round(offset_top + scale * (full_y1 - bb.y1)))
                width_emu = int(round(scale * w_px))
                height_emu = int(round(scale * h_px))
                buf = _io.BytesIO(item["png"]); buf.seek(0)
                slide.shapes.add_picture(buf, left_emu, top_emu, width=width_emu, height=height_emu)
            except Exception:
                continue

    # Save the original full figure using tight layout for comparison
    if save_debug:
        try:
            from .fig_export import save_full_figure_tight
            save_full_figure_tight(fig, f"{(crops_prefix or 'slide')}_full_tight.png", dpi=dpi)
        except Exception:
            pass

    # Save overlays for manual inspection when saving crops or in verbose mode
    if save_debug and save_crops:
        try:
            final_bboxes = [(c["index"], c["bbox_disp"]) for c in crops]
            prefix = crops_prefix or "slide"
            FE.save_full_figure_overview(fig, f"{prefix}_overlay.png", dpi=dpi, axis_bboxes=final_bboxes)
            FE.save_full_figure_component_overlay(fig, f"{prefix}_components.png", dpi=dpi, draw_final_bboxes=final_bboxes)
        except Exception:
            pass

    # Slide numbers: inject an OOXML slide-number field textbox (auto-updating)
    if show_slide_number:
        try:
            from pptx.oxml import parse_xml as _px
            from pptx.oxml.ns import nsdecls as _ns
            from uuid import uuid4 as _uuid4
            # Position box based on desired side
            box_w, box_h = Inches(1.2), Inches(0.35)
            margin = Inches(slide_number_margin_in)
            if str(slide_number_pos).lower() == "left":
                num_left = margin
            else:
                num_left = prs.slide_width - box_w - margin
            num_top = prs.slide_height - box_h - margin
            tb = slide.shapes.add_textbox(num_left, num_top, box_w, box_h)
            tf = tb.text_frame; tf.clear()
            try:
                from pptx.enum.text import MSO_ANCHOR as _ANCH
                tf.margin_left = tf.margin_right = tf.margin_bottom = 0
                tf.vertical_anchor = _ANCH.BOTTOM
            except Exception:
                pass
            # Build field XML; set font size/color; use GUID in braces to match PP expectations
            size = int(slide_number_font_pt * 100)
            color = f"{slide_number_color_rgb[0]:02X}{slide_number_color_rgb[1]:02X}{slide_number_color_rgb[2]:02X}"
            _guid = "{" + str(_uuid4()) + "}"
            fld_xml = (
                f'<a:fld {_ns("a")} id="{_guid}" type="slidenum">'
                f'<a:rPr lang="en-US" smtClean="0" sz="{size}"><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:rPr>'
                '<a:t>1</a:t>'
                '</a:fld>'
            )
            end_rpr = (
                f'<a:endParaRPr {_ns("a")} lang="en-US" sz="{size}"><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:endParaRPr>'
            )
            p = tf.paragraphs[0]._p
            # Remove any default children
            for child in list(p):
                p.remove(child)
            p.append(_px(fld_xml)); p.append(_px(end_rpr))
            try:
                from pptx.enum.text import PP_ALIGN as _AL
                tf.paragraphs[0].alignment = _AL.RIGHT if str(slide_number_pos).lower() != "left" else _AL.LEFT
            except Exception:
                pass
        except Exception:
            pass


def figs_to_deck_precise_combined(
    slides,
    pptx_path,
    template_path: str | None = None,
    *,
    verbose: bool = False,
    save_crops: bool = False,
     save_debug: bool = False,
    crops_dir: str | None = None,
    snap_to_pixel: bool = True,
    write_debug: bool = True,
):
    """Build a single deck with per-slide options.

    slides: list of (fig, kwargs) tuples where kwargs accept the same keys as
    figs_to_deck_precise (except pptx_path).
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pathlib import Path
    from . import fig_export as FE

    # Only use template if explicitly provided
    # Resolve template: explicit path wins; otherwise look in package 'io/templates'
    from pathlib import Path as _Path
    if template_path and not str(template_path).strip():
        template_path = None
    _tpl_note2 = None
    if not template_path:
        _pkg_templates = _Path(__file__).resolve().parent / "templates"
        if _pkg_templates.exists():
            for _pat in ("*.pptx", "*.potx"):
                _cands = [
                    cand
                    for cand in sorted(_pkg_templates.glob(_pat))
                    if not cand.name.startswith("._")
                ]
                if _cands:
                    template_path = str(_cands[0])
                    break
    if template_path and str(template_path).lower().endswith(".potx"):
        _tpl_note2 = f"Ignoring POTX (unsupported by python-pptx): {template_path}"
        template_path = None
    prs = Presentation(template_path) if template_path else Presentation()
    def _clear_template_slides(pres):
        try:
            ids = list(pres.slides._sldIdLst)
            for sldId in ids[::-1]:
                rId = sldId.rId
                pres.slides._sldIdLst.remove(sldId)
                pres.part.drop_rel(rId)
            return len(ids)
        except Exception:
            return 0
    if template_path and len(prs.slides) > 0:
        _clear_template_slides(prs)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Global log header for combined export
    from . import fig_export as FE
    debug_dir = FE._ensure_debug_dir() if save_debug else None
    if debug_dir is not None:
        place_path_global = debug_dir / "pptx_axes_placements.txt"
        with open(place_path_global, "w", encoding="utf-8") as fh:
            print(f"Template path: {template_path if template_path else 'None'}", file=fh)
            print(f"Slide size (EMU): {prs.slide_width} x {prs.slide_height}", file=fh)
            if _tpl_note2:
                print(_tpl_note2, file=fh)

    from pathlib import Path as _P
    if not save_debug:
        if crops_dir is not None and verbose:
            try:
                print("Ignoring crops_dir because save_debug=False")
            except Exception:
                pass
        crops_dir = None
    out_dir_for_crops = _P(crops_dir) if crops_dir else _P(pptx_path).parent
    base_prefix = _P(pptx_path).stem

    # Track SVG sources by slide index so we can attach vector SVG
    # variants to the deck after python-pptx writes the file.  The
    # rasterized PNG produced by the matplotlib pipeline stays as the
    # fallback for older PowerPoint clients.
    slide_idx_to_svg: dict[int, str] = {}
    for idx, entry in enumerate(slides, start=1):
        if isinstance(entry, tuple):
            fig, opts = entry
        else:
            fig = entry.get("fig"); opts = dict(entry)
            opts.pop("fig", None)
        svg_source = getattr(fig, "_ocdkit_svg_source", None)
        if svg_source is not None:
            slide_idx_to_svg[idx] = svg_source
        # compose per-slide crops prefix
        prefix = f"{base_prefix}_s{idx:02d}"
        # Allow per-slide overrides but avoid duplicate kwargs
        slide_save_crops = opts.pop('save_crops', save_crops)
        slide_save_crops = bool(save_debug and slide_save_crops)
        requested_crops_dir = opts.pop('crops_dir', crops_dir)
        slide_crops_dir = requested_crops_dir if save_debug else crops_dir
        slide_crops_prefix = opts.pop('crops_prefix', prefix)
        slide_snap_to_pixel = opts.pop('snap_to_pixel', snap_to_pixel)
        _add_precise_slide(
            prs,
            fig,
            save_crops=slide_save_crops,
            crops_dir=str(out_dir_for_crops if slide_crops_dir is None else slide_crops_dir),
            crops_prefix=slide_crops_prefix,
            snap_to_pixel=slide_snap_to_pixel,
            save_debug=save_debug,
            debug_dir=debug_dir,
            **opts,
        )

    out_path = Path(pptx_path)
    if not write_debug:
        if verbose:
            try:
                print(f"Skipping PPTX export (write_debug=False): {out_path}")
            except Exception:
                pass
        return None
    try:
        out_path.parent.mkdir(parents=True, exist_ok=True)
    except Exception:
        pass
    prs.save(out_path)

    # Attach native SVG variants for any slide whose source was an
    # SvgFigure.  PowerPoint 2016+ renders the vector SVG; older
    # clients use the embedded PNG fallback.
    if slide_idx_to_svg:
        _attach_svg_to_pptx(out_path, slide_idx_to_svg)
        # Make PowerPoint "Convert to Shapes" faithful: correct text size (96dpi
        # bake), no figure distortion (aspect lock), no text shift (zero insets).
        from .figure import apply_convert_to_shapes_fixes
        apply_convert_to_shapes_fixes(out_path)
        if verbose:
            print(f"  attached native SVG to {len(slide_idx_to_svg)} slide(s)")
    if verbose:
        try:
            print(f"Saved PowerPoint: {out_path.resolve()} ({len(prs.slides)} slides)")
            if _tpl_note2:
                print(_tpl_note2)
        except Exception:
            pass




if __name__ == "__main__":
    # Minimal demos: export two slides into debug/ and save PNG crops
    import numpy as np
    import matplotlib.pyplot as plt
    from .fig_export import DPI, _ensure_debug_dir
    from matplotlib.ticker import MaxNLocator

    debug_dir = _ensure_debug_dir()

    # Slide 1: original 2x2 demo
    fig1 = plt.figure(figsize=(10, 6), dpi=DPI)
    x = np.linspace(0, 2 * np.pi, 400)
    gs = fig1.add_gridspec(2, 2, left=0.06, right=0.98, bottom=0.08, top=0.95, wspace=0.3, hspace=0.35)
    ax1 = fig1.add_subplot(gs[0, 0])
    ax1.plot(x, np.sin(x), label="sin"); ax1.plot(x, np.cos(x), label="cos")
    ax1.set_title("Trigs"); ax1.set_xlabel("x"); ax1.set_ylabel("f(x)"); ax1.legend(loc="upper right")
    ax2 = fig1.add_subplot(gs[0, 1])
    rng = np.random.default_rng(0); y = np.sin(3 * x) + 0.2 * rng.standard_normal(len(x))
    ax2.scatter(x, y, s=6, alpha=0.6); ax2.set_title("Noisy sin(3x)"); ax2.set_xlabel("x"); ax2.set_ylabel("y")
    ax3 = fig1.add_subplot(gs[1, 0]); ax3.plot(x, np.tan(x), lw=0.8)
    ax3.set_title("No ticks/text"); ax3.set_xticks([]); ax3.set_yticks([])
    ax4 = fig1.add_subplot(gs[1, 1]); ax4.text(0.5, 0.5, "Center text", ha="center", va="center", fontsize=14)
    ax4.set_title("Center text")
    for ax in fig1.axes:
        ax.yaxis.set_major_locator(MaxNLocator(prune='upper'))
        ax.xaxis.set_major_locator(MaxNLocator(prune='lower'))
        for sp in ax.spines.values():
            sp.set_linewidth(3)
            sp.set_color('gray')
    # Dark theme to match deck
    fig1.patch.set_facecolor("black")
    for ax in fig1.axes:
        ax.set_facecolor("black")
        # Style ticks and labels only when currently visible
        def _color_ticks(axis_obj, tick_color: str, label_color: str) -> None:
            for tick in axis_obj.get_major_ticks():
                if tick.tick1line.get_visible():
                    tick.tick1line.set_color(tick_color)
                if tick.tick2line.get_visible():
                    tick.tick2line.set_color(tick_color)
                if tick.label1.get_visible():
                    tick.label1.set_color(label_color)
                if tick.label2.get_visible():
                    tick.label2.set_color(label_color)
            for tick in axis_obj.get_minor_ticks():
                if tick.tick1line.get_visible():
                    tick.tick1line.set_color(tick_color)
                if tick.tick2line.get_visible():
                    tick.tick2line.set_color(tick_color)
                if tick.label1.get_visible():
                    tick.label1.set_color(label_color)
                if tick.label2.get_visible():
                    tick.label2.set_color(label_color)

        _color_ticks(ax.xaxis, "lime", "white")
        _color_ticks(ax.yaxis, "lime", "white")

        if ax.xaxis.label.get_visible():
            ax.xaxis.label.set_color("white")
        if ax.yaxis.label.get_visible():
            ax.yaxis.label.set_color("white")
        ax.title.set_color("gray")

    # Slide 2: stress-test layout (one wide above two squares)
    fig2 = plt.figure(figsize=(10, 7), dpi=DPI)
    gs2 = fig2.add_gridspec(2, 2, height_ratios=[1.2, 1.0], left=0.07, right=0.98, bottom=0.08, top=0.93, wspace=0.30, hspace=0.42)
    # Top-wide: histogram
    axA = fig2.add_subplot(gs2[0, :])
    rng = np.random.default_rng(42)
    data = rng.normal(loc=0.0, scale=1.0, size=4000)
    axA.hist(data, bins=50, color="#88c", alpha=0.9, edgecolor="#eef")
    axA.set_title("Histogram (top-wide)")
    axA.set_xlabel("value"); axA.set_ylabel("count")
    # Bottom-left: rasterized scatter
    axB = fig2.add_subplot(gs2[1, 0])
    x2 = rng.uniform(-1, 1, 8000)
    y2 = 0.5 * x2 + 0.2 * rng.standard_normal(x2.size)
    axB.scatter(x2, y2, s=4, c="#66c2a5", alpha=0.6, rasterized=True)
    axB.set_title("Rasterized scatter"); axB.set_xlabel("x"); axB.set_ylabel("y")
    # Bottom-right: line plot
    axC = fig2.add_subplot(gs2[1, 1])
    t = np.linspace(0, 10, 600)
    axC.plot(t, np.sin(t), lw=1.5, label="sin")
    axC.plot(t, np.cos(0.6*t), lw=1.2, label="cos")
    axC.legend(loc="upper right", framealpha=0.3)
    axC.set_title("Lines"); axC.set_xlabel("t"); axC.set_ylabel("f(t)")
    # Styling
    for ax in fig2.axes:
        ax.set_facecolor("black")
        for sp in ax.spines.values(): sp.set_linewidth(8); sp.set_color("white")
        ax.tick_params(color="lime", labelcolor="white")
        ax.xaxis.label.set_color("white"); ax.yaxis.label.set_color("white")
        ax.title.set_color("gray")
    fig2.patch.set_facecolor("black")
    # Figure-level annotation between rows (exported as its own textbox)
    fig2.text(0.5, 0.52, "Figure-level note between rows", ha="center", va="center", color="white")

    # Prefer io/templates/*.pptx (template), print what we use
    from pathlib import Path as _Path
    io_templates = _Path(__file__).resolve().parent / "templates"
    chosen_tpl = None
    if io_templates.exists():
        cands = [
            cand for cand in sorted(io_templates.glob("*.pptx")) if not cand.name.startswith("._")
        ]
        if cands:
            chosen_tpl = str(cands[0])

    # Slide 3: spine stress tests (images + selective ticks)
    fig3 = plt.figure(figsize=(10, 7), dpi=DPI)
    gs3 = fig3.add_gridspec(2, 3, width_ratios=[1.1, 1.1, 0.9], left=0.06, right=0.97, bottom=0.08, top=0.93, wspace=0.32, hspace=0.38)

    # 2x2 grid of random noise images with heavy outlines, no ticks/labels
    noise_gs = gs3[:, :2].subgridspec(2, 2, wspace=0.25, hspace=0.25)
    rng = np.random.default_rng(7)
    for idx in range(4):
        ax_img = fig3.add_subplot(noise_gs[idx])
        data = rng.normal(size=(40, 40))
        ax_img.imshow(data, cmap="viridis", origin="lower")
        ax_img.set_xticks([]); ax_img.set_yticks([])
        ax_img.set_facecolor("black")
        for sp in ax_img.spines.values():
            sp.set_linewidth(1)
            sp.set_color("white")

    side_spec = gs3[:, 2].subgridspec(4, 1, hspace=0.24)

    # Line plot with only y ticks (no x ticks)
    ax_y_ticks = fig3.add_subplot(side_spec[0])
    t = np.linspace(0, 4 * np.pi, 400)
    ax_y_ticks.plot(t, np.sin(t), color="#ffaa00", lw=3)
    ax_y_ticks.set_xticks([])
    ax_y_ticks.set_ylabel("sin(t)", color="white")
    ax_y_ticks.tick_params(axis="y", colors="white", width=1)
    for sp in ax_y_ticks.spines.values():
        sp.set_color("white"); sp.set_linewidth(1)
    ax_y_ticks.set_facecolor("black")
    ax_y_ticks.set_title("Y ticks only", color="white")

    # Line plot with only x ticks (no y ticks)
    ax_x_ticks = fig3.add_subplot(side_spec[1])
    ax_x_ticks.plot(t, np.cos(t), color="#55ddff", lw=3)
    ax_x_ticks.set_yticks([])
    ax_x_ticks.set_xlabel("t", color="white")
    ax_x_ticks.tick_params(axis="x", colors="white", width=1)
    for sp in ax_x_ticks.spines.values():
        sp.set_color("white"); sp.set_linewidth(1)
    ax_x_ticks.set_facecolor("black")
    ax_x_ticks.set_title("X ticks only", color="white")

    # Diagonal capstyle stress tests (separate axes)
    ax_diag_butt = fig3.add_subplot(side_spec[2])
    ax_diag_butt.set_facecolor("black")
    ax_diag_butt.set_xticks([]); ax_diag_butt.set_yticks([])
    ax_diag_butt.set_xlim(0, 1); ax_diag_butt.set_ylim(0, 1)
    line_butt, = ax_diag_butt.plot([0.1, 0.9], [0.1, 0.9], color="#ff6699", linewidth=9)
    line_butt.set_solid_capstyle("butt")
    ax_diag_butt.set_title("Diagonal butt cap", color="white", fontsize=12)
    for sp in ax_diag_butt.spines.values():
        sp.set_color("white")
        sp.set_linewidth(1)

    ax_diag_round = fig3.add_subplot(side_spec[3])
    ax_diag_round.set_facecolor("black")
    ax_diag_round.set_xticks([]); ax_diag_round.set_yticks([])
    ax_diag_round.set_xlim(0, 1); ax_diag_round.set_ylim(0, 1)
    line_round, = ax_diag_round.plot([0.1, 0.9], [0.1, 0.9], color="#33ddff", linewidth=9)
    line_round.set_solid_capstyle("round")
    ax_diag_round.set_title("Diagonal round cap", color="white", fontsize=12)
    for sp in ax_diag_round.spines.values():
        sp.set_color("white")
        sp.set_linewidth(1)

    fig3.patch.set_facecolor("black")
    fig3.suptitle("Spine stress grid", color="white", fontsize=18)

    print(f"Using template: {chosen_tpl if chosen_tpl else 'None'}")

    # Single-slide export via combined exporter, saving crops + overlays in debug/
    deck_path = debug_dir / "export_precise_combined.pptx"
    slides_spec = [
        (fig1, dict(
            dpi=DPI,
            background_rgb=(0,0,0),
            title="Centered 2x2 Layout",
            title_font_pt=28,
            title_align="center",
            slide_padding_in=0.3,
            slide_number_pos="right",
            layout_name="Blank",
            save_crops=True,
            crops_dir=str(debug_dir),
        )),
        (fig2, dict(
            dpi=DPI,
            background_rgb=(0,0,0),
            title="Stress Test: Wide + Two Squares",
            title_font_pt=28,
            title_align="center",
            slide_padding_in=0.3,
            slide_number_pos="right",
            layout_name="Blank",
            save_crops=True,
            crops_dir=str(debug_dir),
        )),
        (fig3, dict(
            dpi=DPI,
            background_rgb=(0,0,0),
            title="Spine + Tick Edge Cases",
            title_font_pt=28,
            title_align="center",
            slide_padding_in=0.3,
            slide_number_pos="right",
            layout_name="Blank",
            save_crops=True,
            crops_dir=str(debug_dir),
        )),
    ]
    figs_to_deck_precise_combined(
        slides_spec,
        deck_path,
        template_path=chosen_tpl,
        save_crops=True,
        save_debug=True,
        crops_dir=str(debug_dir),
        verbose=True,
    )

    # Print a short summary and the debug log path
    print(f"Wrote PPTX: {deck_path}")
    dbg_path = debug_dir / "pptx_axes_placements.txt"
    print(f"Placement log: {dbg_path}")
    try:
        # Print the first ~40 lines of the log for quick visibility
        with open(dbg_path, "r", encoding="utf-8") as fh:
            lines = fh.readlines()
        print("--- placement log (head) ---")
        for ln in lines[:40]:
            print(ln.rstrip())
        print("--- end ---")
    except Exception as e:
        print(f"Could not read placement log: {e}")

    # Field tests removed to keep __main__ focused on the single export.

# ── example usage ─────────────────────────────────────────────────────────
# figs = [fig1, fig2, fig3]  # list of previously created figures
# figs_to_deck(figs, os.path.join(base, "plots.pptx"))
